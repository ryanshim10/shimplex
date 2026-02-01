#!/usr/bin/env python3
"""
Shimplex 소개용 PPT 생성 스크립트
실행: python create_ppt.py
출력: Shimplex_Introduction.pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ python-pptx 라이브러리 필요")
    print("설치: pip install python-pptx")
    exit(1)

def create_shimplex_ppt():
    prs = Presentation()
    
    # 색상 정의 (Shimplex 테마: 초록/파랑)
    PRIMARY_GREEN = RGBColor(76, 175, 80)    # #4CAF50
    DARK_GREEN = RGBColor(56, 142, 60)       # #388E3C
    TEXT_DARK = RGBColor(33, 33, 33)         # #212121
    TEXT_GRAY = RGBColor(97, 97, 97)         # #616161
    
    def add_title_slide():
        """슬라이드 1: 타이틀"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
        
        # 메인 타이틀
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Shimplex"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_GREEN
        p.alignment = PP_ALIGN.CENTER
        
        # 서브타이틀
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "개인 AI 클라이언트 / Personal AI Plex"
        p.font.size = Pt(28)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # 설명
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(0.8))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = "어디서든 실행되는 경량 AI 채팅 솔루션"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # 푸터
        footer = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = "github.com/ryanshim10/shimplex | 2026.02"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_problem_slide():
        """슬라이드 2: 문제 제기"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 타이틀
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "왜 Shimplex인가?"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        # 문제점들
        problems = [
            ("🔒 개인정보 유출 우려", "ChatGPT 등 외부 서비스 의존 시 민감 데이터 노출 위험"),
            ("🔧 복잡한 설정", "Docker, CUDA, 의존성 설치... 진입장벽이 너무 높음"),
            ("💻 플랫폼 종속", "Windows용, Mac용 따로 설치해야 하는 번거로움"),
            ("💸 비용 부담", "로컬 AI는 고가 GPU 필요, 클라우드는 월 구독료 발생"),
        ]
        
        y_pos = 1.8
        for icon_title, desc in problems:
            # 아이콘+제목
            box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.5), Inches(0.6))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = icon_title
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_GREEN
            
            # 설명
            desc_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.6), Inches(8), Inches(0.5))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_GRAY
            
            y_pos += 1.2
        
        return slide
    
    def add_solution_slide():
        """슬라이드 3: 솔루션/특징"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 타이틀
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Shimplex 특징"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        # 4가지 특징 (2x2 그리드)
        features = [
            ("🌐", "범용성", "Python만 있으면 OK\nWindows/Mac/Linux 모두 지원", Inches(0.7), Inches(1.8)),
            ("🔌", "유연성", "OpenAI, Claude, Ollama\n모든 주요 LLM 지원", Inches(4.8), Inches(1.8)),
            ("🚀", "간편성", "Docker 없이 실행\n'python app.py' 한 줄로 시작", Inches(0.7), Inches(4.0)),
            ("🛡️", "보안성", "API 키 로컬 관리\n외부 노출 최소화", Inches(4.8), Inches(4.0)),
        ]
        
        for icon, title_text, desc, x, y in features:
            # 카드 배경 (흰색 박스)
            # 제목
            title_box = slide.shapes.add_textbox(x, y, Inches(4), Inches(0.6))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{icon} {title_text}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_GREEN
            
            # 설명
            desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.7), Inches(3.5), Inches(1))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_DARK
            p.line_spacing = 1.3
        
        return slide
    
    def add_architecture_slide():
        """슬라이드 4: 구조 및 시작하기"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 타이틀
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "구조 및 시작하기"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        # 아키텍처 다이어그램 (텍스트로 표현)
        arch_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(2.2))
        tf = arch_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "🏗️  아키텍처"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        
        p = tf.add_paragraph()
        p.text = "사용자 브라우저 → Shimplex 서버(FastAPI) → 외부 LLM API"
        p.font.size = Pt(14)
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "설정: config.json 로컬 파일로 관리"
        p.font.size = Pt(14)
        p.space_before = Pt(4)
        
        # 설치 방법
        install_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.5), Inches(2.5))
        tf = install_box.text_frame
        
        p = tf.paragraphs[0]
        p.text = "⚡  빠른 시작"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        
        p = tf.add_paragraph()
        p.text = "$ git clone https://github.com/ryanshim10/shimplex.git"
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "$ cd shimplex && ./install.sh"
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.space_before = Pt(4)
        
        p = tf.add_paragraph()
        p.text = "$ python app.py"
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.space_before = Pt(4)
        
        # CTA
        cta = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.5), Inches(0.8))
        tf = cta.text_frame
        p = tf.paragraphs[0]
        p.text = "🎯  github.com/ryanshim10/shimplex"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_GREEN
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # 4개 슬라이드 생성
    add_title_slide()
    add_problem_slide()
    add_solution_slide()
    add_architecture_slide()
    
    # 저장
    output_file = "Shimplex_Introduction.pptx"
    prs.save(output_file)
    print(f"✅ PPT 생성 완료: {output_file}")
    print(f"📊 총 {len(prs.slides)}개 슬라이드")
    print("\n슬라이드 목록:")
    print("  1. 타이틀 - Shimplex 소개")
    print("  2. 문제 제기 - 왜 Shimplex인가?")
    print("  3. 특징 - 4가지 핵심 기능")
    print("  4. 구조 및 시작하기 - 아키텍처와 설치법")

if __name__ == "__main__":
    create_shimplex_ppt()
